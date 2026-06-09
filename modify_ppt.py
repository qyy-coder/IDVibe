"""
修改PPT项目计划书，融入技术创新方案
基于: 技术创新方案-详细设计.md
"""
import copy
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# 颜色常量
PRIMARY_BLUE = RGBColor(0x00, 0x7B, 0xFF)
DARK_TEXT = RGBColor(0x34, 0x3A, 0x40)
GRAY_TEXT = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GREEN = RGBColor(0x28, 0xA7, 0x45)
ACCENT_ORANGE = RGBColor(0xFD, 0x7E, 0x14)

SRC = 'AI智能证件照生成与优化系统项目计划书.pptx'
DST = 'AI智能证件照生成与优化系统项目计划书-技术优化版.pptx'

prs = Presentation(SRC)

# 获取幻灯片尺寸
W = prs.slide_width   # 12192000 EMU ≈ 33.867 cm ≈ 13.333 inches
H = prs.slide_height  # 6858000 EMU ≈ 19.05 cm ≈ 7.5 inches

def add_text_box(slide, left, top, width, height, text, font_size=Pt(18),
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text_box(slide, left, top, width, height, paragraphs_data):
    """
    添加富文本文本框
    paragraphs_data: list of list of (text, font_size, bold, color) tuples
    每个元素是一个段落
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, para_runs in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        for j, (text, font_size, bold, color) in enumerate(para_runs):
            if j == 0:
                run = p.add_run()
                run.text = text
                run.font.size = font_size
                run.font.bold = bold
                run.font.color.rgb = color
                run.font.name = 'Microsoft YaHei'
            else:
                run = p.add_run()
                run.text = text
                run.font.size = font_size
                run.font.bold = bold
                run.font.color.rgb = color
                run.font.name = 'Microsoft YaHei'

    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=LIGHT_BG,
                     border_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def clear_slide_shapes(slide):
    """清除幻灯片上除背景外的所有形状"""
    shapes_to_remove = []
    for shape in slide.shapes:
        # 保留背景图片
        if not (shape.shape_type == 13 and shape.name.startswith('Picture')):
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def clone_shape_properties(source_shape, target_shape):
    """复制形状属性(简化版)"""
    try:
        target_shape.left = source_shape.left
        target_shape.top = source_shape.top
        target_shape.width = source_shape.width
        target_shape.height = source_shape.height
    except:
        pass


# ===================================================================
# 1. 修改 Slide 1 (封面) - 更新副标题
# ===================================================================
print("Modifying Slide 1: Title page...")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if hasattr(shape, 'text') and 'HivisionIDPhotos' in shape.text:
        # 保留原有格式，修改文本
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'HivisionIDPhotos' in run.text:
                    run.text = run.text.replace(
                        '基于HivisionIDPhotos开源架构',
                        '自研合规检测引擎 · 本地隐私AI · 多标准智能适配'
                    )
                    run.font.size = Pt(16)
                    run.font.color.rgb = RGBColor(0xDC, 0xE1, 0xE6)
                elif '开源' in run.text:
                    run.text = '八大自研创新技术 · 全流程本地化处理'

# ===================================================================
# 2. 修改 Slide 2 (目录) - 更新章节
# ===================================================================
print("Modifying Slide 2: TOC...")
slide2 = prs.slides[1]
# 找到目录项并修改
toc_items = []
for shape in slide2.shapes:
    if hasattr(shape, 'text') and shape.text.strip():
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text.strip()
                if txt.startswith('0') and '.' in txt:
                    toc_items.append((shape, para, run))

# 更新目录项 (我们需要修改的项目)
toc_updates = {
    '03': ('03. 技术创新亮点   Technical Innovations', True),
    '04': ('04. 产品功能介绍   Product Features', True),
}

for shape, para, run in toc_items:
    txt = run.text.strip()
    if txt.startswith('03.'):
        run.text = '03. 技术创新亮点   Technical Innovations'
    elif txt.startswith('04.'):
        run.text = '04. 产品功能介绍   Product Features'
    elif txt.startswith('05.'):
        run.text = '05. 营销策略           Marketing Strategy'
    elif txt.startswith('06.'):
        run.text = '06. SWOT分析         SWOT Analysis'
    elif txt.startswith('07.'):
        run.text = '07. 风险分析及应对   Risk Analysis & Response'
    elif txt.startswith('08.'):
        run.text = '08. 未来发展计划     Future Plan'


# ===================================================================
# 3. 修改 Slide 3 (项目概述) - 更新核心技术描述
# ===================================================================
print("Modifying Slide 3: Project Overview...")
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if hasattr(shape, 'text') and '以开源架构' in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace(
                    '基于开源架构的AI',
                    '融合八大自研创新技术的AI'
                )
                run.text = run.text.replace(
                    '全流程本地处理',
                    '全流程本地处理 + 智能合规检测'
                )
    if hasattr(shape, 'text') and '以CPU推理' in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace(
                    '以CPU推理为核心',
                    'CPU推理为核心 + 多级模型量化'
                )
                run.text = run.text.replace(
                    '处理响应时间在1秒内',
                    '全流程<1.5秒，模型体积<80MB'
                )
                run.text = run.text.replace(
                    '兼顾速度与效率的同时降低硬件成本',
                    '支持PC/手机/树莓派等多平台部署'
                )


# ===================================================================
# 4. 重做 Slide 4 (核心技术优势 → 技术创新全景)
# ===================================================================
print("Rebuilding Slide 4: Innovation Overview...")
slide4 = prs.slides[3]
clear_slide_shapes(slide4)

# 标题
add_text_box(slide4, Emu(762000), Emu(350000), Emu(10500000), Emu(600000),
             '八大技术创新全景', Pt(36), True, PRIMARY_BLUE)

add_text_box(slide4, Emu(762000), Emu(850000), Emu(10500000), Emu(350000),
             '在开源架构基础上，团队自研8项核心技术，形成完整的技术壁垒',
             Pt(14), False, GRAY_TEXT)

# 8个创新点卡片 (2行×4列)
innovations = [
    ('🔍', '合规检测引擎', '25+项自动检测\nISO/ICAO标准\n全行业首创'),
    ('✂️', '级联抠图流水线', 'Guided Filter边缘优化\n发丝通过率>92%\n自适应回退机制'),
    ('✨', '隐私保护美颜', '本地GAN推理\nArcFace身份保持\n模型仅5MB'),
    ('💡', '光照一致性融合', '球谐光照估计\n软阴影合成\n消除贴纸感'),
    ('🎨', '肤色自适应', 'Fitzpatrick I-VI型\n肤色感知白平衡\n普惠包容设计'),
    ('📐', '拍摄实时引导', '姿态→表情→光照\nAR叠加层反馈\n业界首创'),
    ('⚡', '多级模型量化', '蒸馏→INT8→优化\n树莓派可运行\n总模型<80MB'),
    ('👔', '虚拟正装换装', '轻量扩散模型\n全本地推理\n面部100%不变'),
]

card_w = Emu(2500000)
card_h = Emu(2350000)
start_x = Emu(550000)
start_y = Emu(1300000)
gap_x = Emu(220000)
gap_y = Emu(180000)

for idx, (icon, title, desc) in enumerate(innovations):
    row = idx // 4
    col = idx % 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    # 卡片背景
    card = add_rounded_rect(slide4, x, y, card_w, card_h,
                            fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                            border_color=RGBColor(0xDE, 0xE2, 0xE6))
    # 图标和标题
    add_text_box(slide4, x + Emu(150000), y + Emu(120000),
                 card_w - Emu(300000), Emu(380000),
                 f'{icon}  {title}', Pt(15), True, PRIMARY_BLUE)
    # 描述
    add_text_box(slide4, x + Emu(150000), y + Emu(500000),
                 card_w - Emu(300000), card_h - Emu(650000),
                 desc, Pt(11), False, DARK_TEXT)


# ===================================================================
# 5. 新增 Slide 5 (合规检测引擎详解)
# ===================================================================
print("Adding new slides for key innovations...")

# 新增幻灯片统一追加到末尾（在PowerPoint中手动调整顺序）
# 现有幻灯片索引(0-19)保持不变
def append_new_slide(prs):
    """在末尾添加新幻灯片"""
    blank_layout = prs.slide_layouts[6]  # blank layout
    new_slide = prs.slides.add_slide(blank_layout)
    return new_slide


# 创建合规检测引擎详解页
slide_5 = append_new_slide(prs)

# 标题
add_text_box(slide_5, Emu(762000), Emu(300000), Emu(10500000), Emu(600000),
             '🔍 核心创新：多标准智能合规检测引擎', Pt(32), True, PRIMARY_BLUE)
add_text_box(slide_5, Emu(762000), Emu(800000), Emu(10500000), Emu(300000),
             '全行业首创 — 自动判断证件照是否会被官方审核通过',
             Pt(13), False, GRAY_TEXT)

# 左侧: 检测维度说明
add_text_box(slide_5, Emu(762000), Emu(1200000), Emu(5500000), Emu(350000),
             '25+项智能检测规则', Pt(18), True, DARK_TEXT)

checks_text = (
    '📏 几何检测 (5项)\n'
    '   头部占比70-80% | 眼睛水平线 | 人脸居中 | 下颚位置 | 双肩可见\n\n'
    '🔄 姿态检测 (4项)\n'
    '   偏航角<5° | 俯仰角<5° | 翻滚角<5° | 综合姿态评分\n\n'
    '😊 面部状态检测 (7项)\n'
    '   眼睛开合 | 嘴部闭合 | 中性表情 | 眼镜反光 | 红眼 | 面部遮挡\n\n'
    '💡 光照与色彩检测 (5项)\n'
    '   面部均匀度 | 阴影检测 | 背景均匀度 | 背景色合规 | 整体对比度\n\n'
    '📊 图像质量检测 (5项)\n'
    '   分辨率≥300DPI | 清晰度 | 摩尔纹 | 压缩质量 | 色彩空间sRGB'
)

txBox = slide_5.shapes.add_textbox(Emu(762000), Emu(1600000), Emu(5500000), Emu(4800000))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(checks_text.split('\n')):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = 'Microsoft YaHei'
    if line.startswith('📏') or line.startswith('🔄') or line.startswith('😊') or \
       line.startswith('💡') or line.startswith('📊'):
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = PRIMARY_BLUE
    else:
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = GRAY_TEXT

# 右侧: 对比优势
add_text_box(slide_5, Emu(6700000), Emu(1200000), Emu(5000000), Emu(350000),
             '技术核心：EPnP三维姿态解算', Pt(18), True, DARK_TEXT)

algo_desc = (
    '基于MediaPipe 468点人脸关键点 + 通用3D人脸模型\n'
    '通过solvePnP实时解算Yaw/Pitch/Roll三维姿态角\n'
    '精度达到±2°以内，满足ISO/IEC 19794-5标准要求\n\n'
    '背景均匀度检测采用K-means聚类 + Lab色彩空间\n'
    '8×8网格分块分析，精确识别阴影和色差区域'
)
add_text_box(slide_5, Emu(6700000), Emu(1600000), Emu(5000000), Emu(2000000),
             algo_desc, Pt(12), False, DARK_TEXT)

# 对比表格
add_text_box(slide_5, Emu(6700000), Emu(3300000), Emu(5000000), Emu(350000),
             'vs 竞品合规检测能力', Pt(18), True, DARK_TEXT)

table_text = (
    '                    合规检测  拍摄引导  多国标准\n'
    '本系统              ✅          ✅         ✅\n'
    'HivisionIDPhotos   ❌          ❌         ❌\n'
    '最美证件照          ❌          部分      ❌\n'
    '支付宝证件照        ❌          ❌         ❌\n'
    '传统照相馆          ❌          ❌         ❌'
)
add_text_box(slide_5, Emu(6700000), Emu(3700000), Emu(5000000), Emu(1200000),
             table_text, Pt(11), False, DARK_TEXT)

# 底部: 创新标签
add_rounded_rect(slide_5, Emu(6700000), Emu(5000000), Emu(2300000), Emu(400000),
                 fill_color=PRIMARY_BLUE)
add_text_box(slide_5, Emu(6850000), Emu(5050000), Emu(2000000), Emu(350000),
             '⭐ 全行业首创功能', Pt(14), True, WHITE)


# ===================================================================
# 6. 新增 Slide (技术架构总览)
# ===================================================================
slide_6 = append_new_slide(prs)

add_text_box(slide_6, Emu(762000), Emu(300000), Emu(10500000), Emu(600000),
             '🏗️ 系统技术架构（五层设计）', Pt(32), True, PRIMARY_BLUE)

# 五层架构
layers = [
    ('应用层', '微信小程序  |  Web PWA  |  App (React Native)  |  REST API服务',
     RGBColor(0xE3, 0xF2, 0xFD), PRIMARY_BLUE),
    ('业务编排层', '流水线调度引擎  |  合规规则引擎  |  多国模板管理  |  用户管理系统',
     RGBColor(0xE8, 0xF5, 0xE9), ACCENT_GREEN),
    ('★ AI推理层（11个自研模块）', '人脸检测 | 人像抠图 | 合规检测 | 隐私美颜 | 光照融合 | 肤色校准 | 换装生成 | 质量评分',
     RGBColor(0xFF, 0xF3, 0xE0), ACCENT_ORANGE),
    ('推理引擎层', 'ONNX Runtime  |  OpenCV  |  NumPy  |  CoreML/NNAPI适配',
     RGBColor(0xF3, 0xE5, 0xF5), RGBColor(0x9C, 0x27, 0xB0)),
    ('硬件层', 'CPU (x86/ARM)  |  GPU (可选加速)  |  NPU (移动端神经网络)',
     RGBColor(0xEC, 0xEF, 0xF1), RGBColor(0x54, 0x6E, 0x7A)),
]

y_pos = Emu(1000000)
layer_h = Emu(950000)
for name, content, bg_color, text_color in layers:
    # 层级背景
    rect = add_rounded_rect(slide_6, Emu(762000), y_pos, Emu(10600000), layer_h,
                            fill_color=bg_color)
    # 层级名称
    add_text_box(slide_6, Emu(1000000), y_pos + Emu(80000),
                 Emu(2200000), Emu(350000),
                 name, Pt(16), True, text_color)
    # 内容
    add_text_box(slide_6, Emu(1000000), y_pos + Emu(420000),
                 Emu(10000000), Emu(400000),
                 content, Pt(12), False, DARK_TEXT)
    y_pos += Emu(1050000)

# 底部注释
add_text_box(slide_6, Emu(762000), y_pos + Emu(50000), Emu(10600000), Emu(300000),
             '★ = 团队自研模块（共11个核心模块，基于开源ONNX Runtime + OpenCV推理底座）',
             Pt(11), False, GRAY_TEXT)


# ===================================================================
# 7. 新增 Slide (自研 vs 开源模块边界)
# ===================================================================
slide_7 = append_new_slide(prs)

add_text_box(slide_7, Emu(762000), Emu(300000), Emu(10500000), Emu(600000),
             '📋 自研创新模块 vs 开源依赖（边界清晰）', Pt(28), True, PRIMARY_BLUE)
add_text_box(slide_7, Emu(762000), Emu(750000), Emu(10500000), Emu(300000),
             '评委关注重点：明确标注技术来源与团队贡献',
             Pt(13), False, ACCENT_ORANGE)

# 左侧: 开源依赖
add_text_box(slide_7, Emu(762000), Emu(1100000), Emu(5000000), Emu(400000),
             '📦 依赖的开源组件', Pt(20), True, GRAY_TEXT)

oss_items = [
    ('MTCNN / SCRFD', '人脸检测 (ONNX格式)', 'MIT'),
    ('MODNet', '基础人像抠图', 'Apache 2.0'),
    ('ONNX Runtime', '跨平台推理引擎', 'MIT'),
    ('OpenCV', '图像处理基础库', 'Apache 2.0'),
    ('MediaPipe', '468点人脸关键点', 'Apache 2.0'),
    ('Gradio / FastAPI', 'Web演示与API服务', 'Apache 2.0'),
]

for i, (name, desc, license_type) in enumerate(oss_items):
    y = Emu(1550000) + i * Emu(700000)
    add_rounded_rect(slide_7, Emu(762000), y, Emu(5000000), Emu(600000),
                     fill_color=RGBColor(0xF8, 0xF9, 0xFA),
                     border_color=RGBColor(0xDE, 0xE2, 0xE6))
    add_text_box(slide_7, Emu(1000000), y + Emu(50000), Emu(3000000), Emu(250000),
                 name, Pt(13), True, DARK_TEXT)
    add_text_box(slide_7, Emu(1000000), y + Emu(280000), Emu(4600000), Emu(250000),
                 f'{desc}  |  {license_type}许可证', Pt(10), False, GRAY_TEXT)

# 右侧: 自研模块
add_text_box(slide_7, Emu(6200000), Emu(1100000), Emu(5500000), Emu(400000),
             '🔬 团队自研创新模块 (11个)', Pt(20), True, PRIMARY_BLUE)

self_dev = [
    ('合规检测引擎 ★★★★★', '25+规则 | EPnP姿态解算 | 多国标准'),
    ('级联抠图流水线 ★★★★', 'Guided Filter | 自适应回退 | 发丝优化'),
    ('LQ-BeautyNet ★★★★★', '轻量GAN | ArcFace身份保持 | 5MB部署'),
    ('光照一致性融合 ★★★★', '球谐估计 | 软阴影 | 色彩去溢出'),
    ('肤色自适应管线 ★★★', 'Fitzpatrick分类 | 肤色感知白平衡'),
    ('拍摄实时引导 ★★★★', '3级状态机 | AR叠加层 | 反向引导'),
    ('虚拟正装换装 ★★★★★', 'TinySD蒸馏 | 局部生成 | 面部锁定'),
    ('多级模型量化 ★★★', '蒸馏→INT8→图优化→运行时适配'),
]

for i, (name, desc) in enumerate(self_dev):
    y = Emu(1550000) + i * Emu(575000)
    add_rounded_rect(slide_7, Emu(6200000), y, Emu(5500000), Emu(500000),
                     fill_color=RGBColor(0xE8, 0xF0, 0xFE),
                     border_color=PRIMARY_BLUE)
    add_text_box(slide_7, Emu(6550000), y + Emu(30000), Emu(5000000), Emu(250000),
                 name, Pt(12), True, PRIMARY_BLUE)
    add_text_box(slide_7, Emu(6550000), y + Emu(240000), Emu(5000000), Emu(220000),
                 desc, Pt(10), False, GRAY_TEXT)

# 底部声明
add_text_box(slide_7, Emu(762000), Emu(6350000), Emu(10500000), Emu(350000),
             '核心原则：所有自研模块均100%本地推理，原始照片绝不上传服务器 — 隐私零泄露',
             Pt(12), True, ACCENT_ORANGE, PP_ALIGN.CENTER)


# ===================================================================
# 8. 修改原 Slide 9 (核心创新 → 改为关键算法对比)
# ===================================================================
print("Modifying Slide 9: Algorithm comparison...")
# Slide 9 在原PPT中是第9张 (index 8)
slide9_orig = prs.slides[8]
clear_slide_shapes(slide9_orig)

add_text_box(slide9_orig, Emu(762000), Emu(350000), Emu(10500000), Emu(600000),
             '⚡ 关键技术突破：从能用到好用', Pt(32), True, PRIMARY_BLUE)
add_text_box(slide9_orig, Emu(762000), Emu(850000), Emu(10500000), Emu(300000),
             '基于开源架构的三大维度深度优化，性能与体验质的飞跃',
             Pt(13), False, GRAY_TEXT)

# 三个对比卡片
improvements = [
    ('抠图质量提升', 'MODNet基础抠图\n发丝模糊/锯齿\n半透明区域失真\n通过率~75%',
     '级联精细化流水线\nGuided Filter边缘保持\n自适应置信度回退\n通过率>92%', '发丝通过率\n+17%'),
    ('合规检测能力', '无合规检测\n用户不知照片\n是否会被拒',
     '25+项智能检测\n实时PASS/FAIL反馈\n具体修正建议指导\n全行业首创', '照片一次通过率\n从60%→95%'),
    ('模型部署效率', 'FP32全精度\n~42MB模型\n仅x86 CPU优化',
     'INT8量化部署\n<80MB总大小\nPC/手机/树莓派全平台\nSnapdragon 480可运行', '低端设备\n流畅运行'),
]

for idx, (title, before, after, metric) in enumerate(improvements):
    x = Emu(500000) + idx * Emu(3800000)

    # 标题
    add_text_box(slide9_orig, x, Emu(1250000), Emu(3600000), Emu(400000),
                 title, Pt(18), True, PRIMARY_BLUE, PP_ALIGN.CENTER)

    # Before
    add_rounded_rect(slide9_orig, x, Emu(1750000), Emu(1750000), Emu(2800000),
                     fill_color=RGBColor(0xFF, 0xEB, 0xEE),
                     border_color=RGBColor(0xEF, 0x9A, 0x9A))
    add_text_box(slide9_orig, x + Emu(100000), Emu(1800000), Emu(1600000), Emu(300000),
                 '❌ 优化前', Pt(13), True, RGBColor(0xC6, 0x28, 0x28))
    add_text_box(slide9_orig, x + Emu(100000), Emu(2150000), Emu(1600000), Emu(2200000),
                 before, Pt(11), False, DARK_TEXT)

    # 箭头
    add_text_box(slide9_orig, x + Emu(1800000), Emu(2300000), Emu(300000), Emu(400000),
                 '→', Pt(28), True, PRIMARY_BLUE, PP_ALIGN.CENTER)

    # After
    add_rounded_rect(slide9_orig, x + Emu(1950000), Emu(1750000), Emu(1750000), Emu(2800000),
                     fill_color=RGBColor(0xE8, 0xF5, 0xE9),
                     border_color=RGBColor(0xA5, 0xD6, 0xA7))
    add_text_box(slide9_orig, x + Emu(2050000), Emu(1800000), Emu(1600000), Emu(300000),
                 '✅ 优化后', Pt(13), True, ACCENT_GREEN)
    add_text_box(slide9_orig, x + Emu(2050000), Emu(2150000), Emu(1600000), Emu(2200000),
                 after, Pt(11), False, DARK_TEXT)

    # 底部指标
    add_rounded_rect(slide9_orig, x + Emu(500000), Emu(4700000), Emu(2800000), Emu(380000),
                     fill_color=PRIMARY_BLUE)
    add_text_box(slide9_orig, x + Emu(600000), Emu(4750000), Emu(2600000), Emu(350000),
                 f'📈 {metric}', Pt(13), True, WHITE, PP_ALIGN.CENTER)


# ===================================================================
# 9. 修改原 Slide 18 (未来发展计划) - index 17
# ===================================================================
print("Modifying Slide 18: Future Plans...")
slide18 = prs.slides[17]
# 寻找包含年份的shape并更新
for shape in slide18.shapes:
    if hasattr(shape, 'text'):
        txt = shape.text
        if '2026年内' in txt or '短期' in txt:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = run.text.replace('完成微信小程序上线及封闭测试',
                                                '完成合规检测引擎v1.0 + 微信小程序上线')
                    run.text = run.text.replace('优化无障碍功能体验',
                                                '级联抠图 + 光照融合 + 拍摄引导上线')
                    run.text = run.text.replace('冲刺省级大赛奖项', '冲刺省赛金奖，申请软著2项')
        elif '2027年' in txt or '中期' in txt:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = run.text.replace('推出App多端版本及边缘适配',
                                                '推出iOS/Android + 合规检测引擎v2.0 (25项全规则)')
                    run.text = run.text.replace('付费转化率提升至15%',
                                                'LQ-BeautyNet + 虚拟换装上线，付费转化>15%')
                    run.text = run.text.replace('用户规模突破30万', '用户规模突破30万，B端API签约5家企业')
        elif '2028年' in txt or '长期' in txt:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = run.text.replace('对接政务与企业API，规模化应用',
                                                '对接政务/HR/考试报名系统API，规模化应用')
                    run.text = run.text.replace('推动隐私AI证件照成行业标准',
                                                '推动隐私AI合规检测成为证件照行业标准')


# ===================================================================
# 10. 保存
# ===================================================================
print(f"\nSaving to: {DST}")
prs.save(DST)
print("Done! Modified PPT saved successfully.")
print(f"\nModified slides summary:")
print(f"  Slide 1: Updated title subtitle")
print(f"  Slide 2: Updated TOC structure")
print(f"  Slide 3: Updated project overview text")
print(f"  Slide 4: Replaced with 8-Innovation Overview cards")
print(f"  Slide 5-9: (Originals - Pain Points, Market, etc.)")
print(f"  Slide 10: Replaced with Algorithm Improvements Comparison")
print(f"  Slides 11-19: Original remaining slides")
print(f"  Slide 20: Updated Future Plan timeline")
print(f"  Slide 21: [NEW] Compliance Detection Engine detail")
print(f"  Slide 22: [NEW] System Architecture (5-layer)")
print(f"  Slide 23: [NEW] Self-developed vs Open-source boundary")
print(f"\nTotal slides: {len(prs.slides)} (original 20 + 3 new)")
print(f"NOTE: New slides are appended at the end.")
print(f"      Please manually reorder in PowerPoint:")
print(f"      Move slide 21-23 → between slide 4 and slide 5")
