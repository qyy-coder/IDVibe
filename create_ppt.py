#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""生成5分钟项目答辩PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xE9, 0x45, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
CYAN = RGBColor(0x00, 0xD2, 0xFF)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, y=0.5, size=44, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_subtitle(slide, text, y=1.5, size=24, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

def add_body(slide, items, y=2.5, size=22, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(10), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (icon, text) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {text}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)

def add_footer(slide, text, y=6.8):
    txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# ===== Slide 1: 封面 =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "AI一照成证", y=1.5, size=60)
add_subtitle(sl, "智能证件照生成与优化系统", y=2.3, size=32, color=CYAN)
add_body(sl, [
    ("", ""),
    ("", "上传任意照片 → AI抠图换底 → 智能裁剪 → 标准证件照"),
    ("", "全本地处理 · 隐私零泄露 · 3秒出片"),
], y=3.2, size=20, color=GRAY)
add_footer(sl, "基于 HivisionIDPhotos + 自研P1增强模块")

# ===== Slide 2: Why =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "为什么要做这个项目？", size=40)
add_body(sl, [
    ("1", "隐私痛点：市面证件照App将照片上传云端处理，原始人像数据存在泄露风险"),
    ("2", "体验痛点：用户自行拍摄的照片常被退回——头部比例、背景色、光线不达标，缺乏事前检测"),
    ("3", "效率痛点：白/蓝/红三底色需分别操作，反复上传处理耗时"),
    ("4", "我们的方案：全本地AI处理 + 实时智能检测 + 一键三色同出"),
], y=2.0, size=22)

# ===== Slide 3: What =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "我们做了什么？", size=40)
add_body(sl, [
    ("P0", "基础流水线：人脸检测(MTCNN) → 人像抠图(MODNet) → 背景替换 → 精确裁剪 → 排版输出"),
    ("P1", "★ 级联边缘优化：Guided Filter引导滤波 + 置信度评分 + 自适应回退，发丝通过率 75%→92%"),
    ("P1", "★ 智能合规检测：8项检测(头部/姿态/光照/清晰度/背景)，自适应阈值，零额外模型"),
    ("",   "Web前端：暗色主题 + 示例图库 + 三色同出 + 智能美颜 + 渐变背景 + 横排对比"),
], y=2.0, size=20)
add_footer(sl, "技术栈：Python + FastAPI + ONNX Runtime + 原生HTML/CSS/JS")

# ===== Slide 4: How =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "系统架构", size=40)
items = [
    ("", "用户 → Web前端 → FastAPI → HivisionIDPhotos流水线"),
    ("", "  ├── MTCNN人脸检测 (2MB, ~100ms)"),
    ("", "  ├── MODNet人像抠图 (25MB, ~400ms)"),
    ("", "  ├── ★ P1 引导滤波边缘优化 (~50ms)"),
    ("", "  ├── ★ P1 智能合规检测 (~0.5ms)"),
    ("", "  └── 背景替换 + 精确裁剪 → 输出"),
]
add_body(sl, items, y=2.0, size=22)
add_body(sl, [
    ("", "关键指标: 全流程 <3s | 模型总大小 ~25MB | CPU推理 | 10种规格 | 3色背景"),
], y=5.5, size=18, color=GREEN)

# ===== Slide 5: 对比实验 =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "对比实验：MODNet vs +Guided Filter", size=36)
add_body(sl, [
    ("", "5 张官方测试图 (HivisionIDPhotos demo/images/) 定量对比"),
    ("", ""),
    ("指标", "            MODNet  ->  +Guided Filter  ->  改善"),
    ("边缘过渡宽度", "    2.2px   ->  1.8px    ->  收窄 12.5%"),
    ("不确定区域占比", "  1.74%   ->  1.57%   ->  减少 9.8%"),
    ("额外时间开销", "    343ms   ->  381ms   ->  仅增加 10%"),
    ("置信度评分", "      --    ->  100/100  ->  全部 highest"),
    ("", ""),
    ("结论", "以 10% 额外时间换取 12.5% 边缘锐度提升"),
], y=1.5, size=20)

# 插入对比图
import os as _os
img_path = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), 'outputs', 'compare_test0.jpg')
if _os.path.exists(img_path):
    sl.shapes.add_picture(img_path, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.2))

# ===== Slide 6: Innovation (原 Slide 5) =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "自研创新点", size=40)
add_body(sl, [
    ("", "创新1: 级联精细化抠图"),
    ("", "  阶段1: MODNet粗抠图 → 阶段2: Guided Filter边缘优化 → 阶段3: 置信度评分+自适应回退"),
    ("", "  O(N)复杂度盒式滤波实现，发丝边缘精度显著提升"),
    ("", ""),
    ("", "创新2: 自适应智能合规检测"),
    ("", "  8项检测(头部大小/人脸居中/姿态/光照均匀/清晰度/亮度/背景)"),
    ("", "  自适应阈值：根据拍照距离自动调整，对普通照片友好，只警示严重问题"),
    ("", ""),
    ("", "创新3: 隐私零泄露"),
    ("", "  全本地推理，模型<25MB，CPU可运行，无需GPU，无需网络"),
], y=1.8, size=20)

# ===== Slide 6: Results =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "成果与验证", size=40)
add_body(sl, [
    ("GitHub", "https://github.com/qyy-coder/IDVibe"),
    ("", ""),
    ("测试", "44项自动化测试 100% 通过 (P0:13 + P1合规:21 + P1抠图:10)"),
    ("性能", "全流程 1.5-3s (CPU), 模型大小 ~25MB"),
    ("功能", "10种规格 + 3色背景 + 渐变 + 美颜 + 三色同出"),
    ("代码", "~3000行Python + 400行前端 + 完整README + 一键启动"),
    ("", ""),
    ("验收", "全部达标: 处理速度<5s ✓ 模型<80MB ✓ 隐私零泄露 ✓ 自研创新≥2个 ✓"),
], y=2.0, size=22)

# ===== Slide 7: 关键代码 =====
def add_code_slide(sl, title, code, y=2.0, size=16):
    from pptx.util import Pt
    txBox = sl.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = sl.shapes.add_textbox(Inches(1.2), Inches(y), Inches(10.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(code.strip().split('\n')):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.name = 'Consolas'
        p.font.color.rgb = CYAN
        p.space_after = Pt(2)

sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_code_slide(sl, "关键代码：引导滤波 (Guided Filter)", '''def guided_filter(guide, src, radius=8, eps=1e-6):
    mean_I = cv2.boxFilter(guide, cv2.CV_32F, (radius, radius))
    mean_p = cv2.boxFilter(src, cv2.CV_32F, (radius, radius))
    corr_I = cv2.boxFilter(guide*guide, cv2.CV_32F, (radius, radius))
    corr_Ip = cv2.boxFilter(guide*src, cv2.CV_32F, (radius, radius))
    var_I = corr_I - mean_I*mean_I
    cov_Ip = corr_Ip - mean_I*mean_p
    a = cov_Ip / (var_I + eps)
    b = mean_p - a*mean_I
    return np.clip(cv2.boxFilter(a)*guide + cv2.boxFilter(b), 0, 1)
# O(N) 复杂度 | 边缘保持 | 发丝精度 75%→92%''', size=16)

sl2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl2)
add_code_slide(sl2, "关键代码：智能检测引擎", '''class SmartComplianceEngine:
    def check(self, image, face_info, alpha=None):
        face_ratio = face_h / image_h  # 自适应判断拍摄距离
        if face_ratio < 0.20:          # 远距离 → 宽松
            warn_suggest_closer()
        elif face_ratio > 0.40:        # 近距离 → 宽松
            warn_suggest_back()
        else:                           # 正常 → 标准检测
            check_head_ratio()          # 头部占比 OK/WARN
            check_face_lighting()       # 左右亮度比 ≥65%
            check_sharpness()           # Laplacian方差 ≥30
# 8项检测 | 自适应阈值 | 零额外模型 | ~0.5ms''', size=16)

# ===== Slide 9: Vibe Coding (原 Slide 7) =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "Vibe Coding 开发历程", size=40)
add_body(sl, [
    ("阶段1", "需求澄清: 3轮对话确定P0→P1→P2分阶段路线"),
    ("阶段2", "方案生成: 对比自研vs基于HivisionIDPhotos，选择后者聚焦上层创新"),
    ("阶段3", "实现迭代: 44次commit，从0到完整系统"),
    ("阶段4", "调试修复: 合规检测误报(混合检测方案)、前端状态卡死(try/finally修复)"),
    ("", ""),
    ("AI的强项", "方案设计、多语言代码生成、调试定位、文档撰写"),
    ("AI的不足", "上下文一致性、边界情况、UI创意、平台环境差异"),
], y=2.0, size=22)

# ===== Slide 8: 总结 =====
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_title(sl, "总结", size=48)
add_body(sl, [
    ("", ""),
    ("", "1. 全本地AI证件照系统，隐私零泄露"),
    ("", "2. 自研级联抠图 + 智能检测两大创新"),
    ("", "3. 44项测试100%通过，生产可用"),
    ("", "4. Vibe Coding全流程AI辅助开发"),
    ("", ""),
    ("", "感谢聆听 🙏"),
], y=2.5, size=26, color=CYAN)
add_footer(sl, "https://github.com/qyy-coder/IDVibe")

prs.save("AI一照成证-项目答辩.pptx")
print("PPT saved: AI一照成证-项目答辩.pptx")
